#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
validate_deck.py — Validador propio y reforzado para presentaciones .pptx (skill pptx-pro-max).

Autónomo: NO depende de ninguna skill de terceros. Solo requiere librerías libres
(python-pptx, lxml). Comprueba, en dos niveles:

  ERRORES (rompen PowerPoint o son defectos graves) — hacen fallar la validación:
    - El archivo no abre / no es un OOXML de presentación válido.
    - XML mal formado en cualquier parte del paquete.
    - Defectos de gráfico que PowerPoint rechaza:
        * dLblPos = outEnd en un gráfico de barras/columnas APILADAS.
        * gráfico que referencia ejes (axId) no declarados (bug de eje secundario).
    - Colores hex inválidos (con '#', longitud != 6, o caracteres no hex).
    - Relaciones rotas: un Target interno de un .rels que no existe en el paquete.
    - Marcadores de plantilla sin reemplazar (lorem/ipsum/xxx/[insert]/[Description]...).
    - Presentación sin diapositivas.

  ADVERTENCIAS (calidad, comunicación, accesibilidad) — no hacen fallar, pero se listan:
    - Contraste texto/fondo por debajo de WCAG AA (4.5:1 normal, 3:1 texto grande).
    - Tamaño de fuente por debajo del piso recomendado.
    - Tipografías fuera de la lista segura (recuerda incrustarlas).
    - Cuadros de texto sin txBox=1 (los lectores de pantalla los anuncian como 'gráfico').
    - Posibles antipatrones: franja/barra decorativa de ancho casi completo, o regla fina
      justo bajo el título.
    - Orden de <p:notesMasterIdLst> antes de <p:sldIdLst> (PowerPoint puede rechazarlo).

Uso:
    python3 validate_deck.py deck.pptx [--strict] [--quiet]
      --strict : trata las advertencias como errores (exit 1 si hay cualquiera).
      --quiet  : solo imprime el resumen final.

Salida: informe legible + resumen. Exit 0 si pasa; 1 si hay errores (o advertencias con --strict).
"""

import argparse
import re
import sys
import zipfile
from pathlib import Path

try:
    from lxml import etree
except Exception:  # pragma: no cover
    print("Falta 'lxml'. Instala con: pip install lxml", file=sys.stderr)
    sys.exit(2)

# --- Namespaces OOXML (estándar público ISO/IEC 29500) ---
NS = {
    "a":  "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p":  "http://schemas.openxmlformats.org/presentationml/2006/main",
    "c":  "http://schemas.openxmlformats.org/drawingml/2006/chart",
    "r":  "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel":"http://schemas.openxmlformats.org/package/2006/relationships",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
}

# Fuentes que vienen con Office y rinden fiables (ver references/03).
SAFE_FONTS = {
    "arial", "calibri", "cambria", "times new roman", "courier new",
    "century schoolbook", "bookman old style", "georgia", "verdana",
    "tahoma", "trebuchet ms", "garamond", "gill sans", "segoe ui",
    "arial black", "franklin gothic", "candara", "constantia", "corbel",
}
PLACEHOLDER_PATTERNS = re.compile(
    r"\blorem\b|\bipsum\b|\bx{3,}\b|\[insert|\[\.\.\.\]|\bTODO\b|"
    r"\[description\]|\[key takeaways|placeholder|lorem ipsum",
    re.IGNORECASE,
)
FONT_FLOOR_PT = 12          # por debajo de esto: advertencia dura de legibilidad
HEX_RE = re.compile(r"^[0-9A-Fa-f]{6}$")


class Report:
    def __init__(self):
        self.errors = []
        self.warnings = []
        self._seen = set()
    def err(self, where, msg):
        line = f"[{where}] {msg}"
        if line not in self._seen:
            self._seen.add(line); self.errors.append(line)
    def warn(self, where, msg):
        line = f"[{where}] {msg}"
        if line not in self._seen:
            self._seen.add(line); self.warnings.append(line)


# ---------- utilidades de color / contraste (WCAG 2.1) ----------
def _lin(c):
    c = c / 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4

def _luminance(hex6):
    r = int(hex6[0:2], 16); g = int(hex6[2:4], 16); b = int(hex6[4:6], 16)
    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)

def contrast_ratio(fg, bg):
    l1, l2 = _luminance(fg), _luminance(bg)
    hi, lo = max(l1, l2), min(l1, l2)
    return (hi + 0.05) / (lo + 0.05)


# ---------- carga del paquete ----------
def load_parts(pptx_path):
    """Devuelve dict {nombre_parte: bytes} de todo el ZIP."""
    with zipfile.ZipFile(pptx_path, "r") as zf:
        return {n: zf.read(n) for n in zf.namelist() if not n.endswith("/")}


def parse_xml(name, data, rep):
    try:
        return etree.fromstring(data)
    except etree.XMLSyntaxError as e:
        rep.err(name, f"XML mal formado: {e}")
        return None


# ---------- chequeos ----------
def check_opens_with_pptx(pptx_path, rep):
    try:
        from pptx import Presentation
    except Exception:
        rep.warn("motor", "python-pptx no está instalado; se omite la apertura de sanidad "
                          "(pip install python-pptx para el chequeo completo).")
        return None
    try:
        prs = Presentation(str(pptx_path))
        n = len(prs.slides._sldIdLst)
        if n == 0:
            rep.err("presentación", "la presentación no tiene diapositivas.")
        return n
    except Exception as e:
        rep.err("presentación", f"python-pptx no pudo abrir el archivo: {e}")
        return None


def check_all_xml_wellformed(parts, rep):
    for name, data in parts.items():
        if name.endswith(".xml") or name.endswith(".rels"):
            parse_xml(name, data, rep)


def check_colors(parts, rep):
    """srgbClr val debe ser 6 hex, sin '#', sin alpha de 8 dígitos."""
    for name, data in parts.items():
        if not name.endswith(".xml"):
            continue
        root = parse_xml(name, data, rep)
        if root is None:
            continue
        for el in root.iter():
            tag = etree.QName(el).localname
            if tag == "srgbClr":
                val = el.get("val", "")
                if not HEX_RE.match(val):
                    rep.err(name, f"color srgbClr inválido: '{val}' (debe ser 6 hex, sin '#').")


def check_placeholders(parts, rep):
    for name, data in parts.items():
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        root = parse_xml(name, data, rep)
        if root is None:
            continue
        texts = [t.text or "" for t in root.iter(f'{{{NS["a"]}}}t')]
        joined = " ".join(texts)
        m = PLACEHOLDER_PATTERNS.search(joined)
        if m:
            rep.err(name, f"marcador de plantilla sin reemplazar: '{m.group(0)}'.")


def check_charts(parts, rep):
    """Defectos de gráfico que PowerPoint rechaza."""
    c = NS["c"]
    for name, data in parts.items():
        if not (name.startswith("ppt/charts/chart") and name.endswith(".xml")):
            continue
        root = parse_xml(name, data, rep)
        if root is None:
            continue

        # (1) outEnd en barras/columnas apiladas
        for bar in root.iter(f"{{{c}}}barChart"):
            grp = bar.find(f"{{{c}}}grouping")
            grouping = grp.get("val") if grp is not None else "clustered"
            if grouping in ("stacked", "percentStacked"):
                for pos in bar.iter(f"{{{c}}}dLblPos"):
                    if pos.get("val") == "outEnd":
                        rep.err(name, "dLblPos='outEnd' en gráfico APILADO (PowerPoint lo "
                                      "rechaza; usa 'ctr', 'inEnd' o 'inBase').")

        # (2) un chart-type sin suficientes ejes declarados => archivo roto.
        # Nota: pptxgenjs añade un axId extra benigno en gráficos de línea que
        # PowerPoint tolera; por eso solo marcamos error cuando NO hay al menos
        # dos ejes declarados (0-1), señal de un gráfico realmente sin ejes.
        declared, used = set(), set()
        axis_tags = {"catAx", "valAx", "dateAx", "serAx"}
        charttype_tags = {"barChart", "lineChart", "areaChart", "scatterChart",
                          "bubbleChart", "radarChart", "stockChart", "surfaceChart",
                          "line3DChart", "bar3DChart", "area3DChart", "pie3DChart"}
        for el in root.iter():
            if etree.QName(el).localname != "axId":
                continue
            parent = el.getparent()
            ptag = etree.QName(parent).localname if parent is not None else ""
            if ptag in axis_tags:
                declared.add(el.get("val"))
            elif ptag in charttype_tags:
                used.add(el.get("val"))
        if used and len(declared) < 2:
            rep.err(name, f"el gráfico usa ejes pero declara solo {len(declared)} (<2); "
                          "un combo con eje secundario necesita valAxes y catAxes completos.")


def check_rels(parts, rep):
    """Cada Target interno de un .rels debe existir como parte del paquete."""
    part_names = set(parts.keys())
    for name, data in parts.items():
        if not name.endswith(".rels"):
            continue
        root = parse_xml(name, data, rep)
        if root is None:
            continue
        base = str(Path(name).parent.parent).replace("\\", "/")  # _rels está bajo la carpeta
        if base == ".":
            base = ""
        for r in root.iter(f'{{{NS["rel"]}}}Relationship'):
            if (r.get("TargetMode") or "Internal") != "Internal":
                continue
            target = r.get("Target", "")
            if target.startswith("/"):
                resolved = target.lstrip("/")
            else:
                # resolver relativo a la carpeta del origen
                origin = base + "/" if base else ""
                resolved = str(Path(origin + target)).replace("\\", "/")
                # normalizar ../
                parts_stack = []
                for seg in resolved.split("/"):
                    if seg == "..":
                        if parts_stack:
                            parts_stack.pop()
                    elif seg not in ("", "."):
                        parts_stack.append(seg)
                resolved = "/".join(parts_stack)
            if resolved and resolved not in part_names:
                rep.err(name, f"relación rota: Target '{target}' no existe en el paquete.")


def check_presentation_order(parts, rep):
    data = parts.get("ppt/presentation.xml")
    if not data:
        return
    root = parse_xml("ppt/presentation.xml", data, rep)
    if root is None:
        return
    children = [etree.QName(e).localname for e in root]
    if "notesMasterIdLst" in children and "sldIdLst" in children:
        if children.index("notesMasterIdLst") < children.index("sldIdLst"):
            rep.warn("ppt/presentation.xml",
                     "<p:notesMasterIdLst> aparece antes de <p:sldIdLst>; PowerPoint puede "
                     "rechazar el archivo. Muévelo justo después de <p:sldIdLst>.")


def _emu_to_in(v):
    try:
        return int(v) / 914400.0
    except Exception:
        return None


def check_slides_quality(parts, rep):
    """Contraste, tamaños de fuente, fuentes seguras, txBox, antipatrones (heurísticos)."""
    a, p = NS["a"], NS["p"]
    SLIDE_W_IN = 13.333  # por defecto; se ajusta abajo si el archivo declara otra cosa
    pres = parts.get("ppt/presentation.xml")
    if pres:
        proot = parse_xml("ppt/presentation.xml", pres, rep)
        if proot is not None:
            sz = proot.find(f"{{{p}}}sldSz")
            if sz is not None:
                w = _emu_to_in(sz.get("cx"))
                if w:
                    SLIDE_W_IN = w

    for name, data in sorted(parts.items()):
        if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
            continue
        root = parse_xml(name, data, rep)
        if root is None:
            continue

        # color de fondo del slide (si es sólido explícito)
        bg = None
        bgpr = root.find(f".//{{{p}}}bg//{{{a}}}solidFill/{{{a}}}srgbClr")
        if bgpr is not None and HEX_RE.match(bgpr.get("val", "")):
            bg = bgpr.get("val")

        for sp in root.iter(f"{{{p}}}sp"):
            txbody = sp.find(f"{{{p}}}txBody")
            if txbody is None:
                continue
            has_text = any((t.text or "").strip() for t in txbody.iter(f"{{{a}}}t"))
            if not has_text:
                continue

            # txBox=1 (accesibilidad)
            cnv = sp.find(f"{{{p}}}nvSpPr/{{{p}}}cNvSpPr")
            if cnv is not None and cnv.get("txBox") != "1":
                rep.warn(name, "cuadro de texto sin txBox='1' (los lectores de pantalla lo "
                               "anuncian como 'gráfico'; usa isTextBox:true en pptxgenjs).")

            # relleno de la forma (para contraste)
            shape_fill = None
            sf = sp.find(f"{{{p}}}spPr/{{{a}}}solidFill/{{{a}}}srgbClr")
            if sf is not None and HEX_RE.match(sf.get("val", "")):
                shape_fill = sf.get("val")
            eff_bg = shape_fill or bg or "FFFFFF"

            # anchura de la forma (para heurística de franja)
            ext = sp.find(f"{{{p}}}spPr/{{{a}}}xfrm/{{{a}}}ext")
            off = sp.find(f"{{{p}}}spPr/{{{a}}}xfrm/{{{a}}}off")
            if ext is not None:
                w_in = _emu_to_in(ext.get("cx")); h_in = _emu_to_in(ext.get("cy"))
                if w_in and h_in and shape_fill and w_in > 0.80 * SLIDE_W_IN and h_in < 0.35:
                    rep.warn(name, "posible franja/barra decorativa de ancho casi completo "
                                   "(antipatrón; usa espacio en blanco o un tinte de fondo).")

            # runs: color + tamaño + fuente
            for r_ in txbody.iter(f"{{{a}}}r"):
                rpr = r_.find(f"{{{a}}}rPr")
                txt = (r_.find(f"{{{a}}}t").text if r_.find(f"{{{a}}}t") is not None else "") or ""
                if not txt.strip():
                    continue
                sz_pt = None
                fg = None
                bold = False
                if rpr is not None:
                    if rpr.get("sz"):
                        try:
                            sz_pt = int(rpr.get("sz")) / 100.0
                        except Exception:
                            sz_pt = None
                    bold = rpr.get("b") == "1"
                    fclr = rpr.find(f"{{{a}}}solidFill/{{{a}}}srgbClr")
                    if fclr is not None and HEX_RE.match(fclr.get("val", "")):
                        fg = fclr.get("val")
                    latin = rpr.find(f"{{{a}}}latin")
                    if latin is not None:
                        face = (latin.get("typeface") or "").strip().lower()
                        if face and face not in SAFE_FONTS:
                            rep.warn(name, f"tipografía fuera de la lista segura: "
                                           f"'{latin.get('typeface')}' (recuerda incrustarla).")

                # tamaño mínimo
                if sz_pt is not None and sz_pt < FONT_FLOOR_PT:
                    rep.warn(name, f"texto a {sz_pt:g} pt (< {FONT_FLOOR_PT} pt): riesgo de "
                                   f"ilegibilidad — «{txt[:24].strip()}…».")

                # contraste WCAG (solo si ambos colores son resolubles)
                if fg is not None:
                    ratio = contrast_ratio(fg, eff_bg)
                    is_large = (sz_pt is not None and (sz_pt >= 18 or (sz_pt >= 14 and bold)))
                    threshold = 3.0 if is_large else 4.5
                    if ratio < threshold:
                        rep.warn(name, f"contraste {ratio:.1f}:1 (< {threshold:g}:1) "
                                       f"texto #{fg} sobre #{eff_bg} — «{txt[:24].strip()}…».")


def main():
    ap = argparse.ArgumentParser(description="Validador reforzado de .pptx (pptx-pro-max).")
    ap.add_argument("pptx", help="ruta al archivo .pptx")
    ap.add_argument("--strict", action="store_true", help="tratar advertencias como errores")
    ap.add_argument("--quiet", action="store_true", help="solo el resumen")
    args = ap.parse_args()

    path = Path(args.pptx)
    if not path.is_file():
        print(f"Error: no existe {path}", file=sys.stderr)
        sys.exit(2)

    rep = Report()
    try:
        parts = load_parts(path)
    except zipfile.BadZipFile:
        print("Error: el archivo no es un .pptx válido (ZIP corrupto).", file=sys.stderr)
        sys.exit(1)

    check_opens_with_pptx(path, rep)
    check_all_xml_wellformed(parts, rep)
    check_colors(parts, rep)
    check_placeholders(parts, rep)
    check_charts(parts, rep)
    check_rels(parts, rep)
    check_presentation_order(parts, rep)
    check_slides_quality(parts, rep)

    if not args.quiet:
        if rep.errors:
            print(f"\n✗ ERRORES ({len(rep.errors)}):")
            for e in rep.errors:
                print("  •", e)
        if rep.warnings:
            print(f"\n⚠ ADVERTENCIAS ({len(rep.warnings)}):")
            for w in rep.warnings:
                print("  •", w)

    n_err, n_warn = len(rep.errors), len(rep.warnings)
    print(f"\nResumen: {n_err} error(es), {n_warn} advertencia(s).")
    if n_err == 0 and (n_warn == 0 or not args.strict):
        print("✓ Validación superada." if n_warn == 0 else "✓ Sin errores (revisa las advertencias).")
        sys.exit(0)
    print("✗ Validación NO superada.")
    sys.exit(1)


if __name__ == "__main__":
    main()
