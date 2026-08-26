#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
check_deck.py — QA del .pptx de imágenes producido por deck-studio.
Autónomo: solo requiere python-pptx.

Comprueba:
  - el .pptx abre y tiene al menos una lámina;
  - cada lámina lleva exactamente una imagen a pantalla completa (a sangre);
  - si se pasa --png, que el nº de láminas == nº de PNG renderizados.

Uso:
    python3 check_deck.py <deck.pptx> [--png <dir_png>]
Salida: informe + exit 0 si pasa, 1 si hay problemas.
"""
import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except Exception:
    print("Falta python-pptx. Instala: pip install python-pptx", file=sys.stderr)
    sys.exit(2)

FULL_W = 13.333
FULL_H = 7.5
TOL = 0.15  # pulgadas de tolerancia para 'a sangre'


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--png", default=None)
    args = ap.parse_args()

    p = Path(args.pptx)
    if not p.is_file():
        print(f"No existe {p}", file=sys.stderr); sys.exit(2)

    errors = []
    try:
        prs = Presentation(str(p))
    except Exception as e:
        print(f"✗ No se pudo abrir el .pptx: {e}"); sys.exit(1)

    slides = list(prs.slides)
    if not slides:
        errors.append("el .pptx no tiene láminas.")

    sw = prs.slide_width / 914400.0
    sh = prs.slide_height / 914400.0

    for i, sl in enumerate(slides, 1):
        pics = [sh_ for sh_ in sl.shapes if sh_.shape_type == 13]  # 13 = PICTURE
        if len(pics) == 0:
            errors.append(f"lámina {i}: no tiene imagen (deck-studio empaqueta una imagen por lámina).")
            continue
        if len(pics) > 1:
            errors.append(f"lámina {i}: tiene {len(pics)} imágenes (se espera exactamente 1 a sangre).")
        pic = pics[0]
        w = pic.width / 914400.0; h = pic.height / 914400.0
        if abs(w - sw) > TOL or abs(h - sh) > TOL:
            errors.append(f"lámina {i}: la imagen no cubre la lámina (mide {w:.2f}×{h:.2f}\", "
                          f"lámina {sw:.2f}×{sh:.2f}\").")

    if args.png:
        n_png = len(sorted(Path(args.png).glob("*.png")))
        if n_png != len(slides):
            errors.append(f"nº de PNG ({n_png}) ≠ nº de láminas ({len(slides)}).")

    if errors:
        print(f"✗ {len(errors)} problema(s):")
        for e in errors:
            print("  •", e)
        print(f"\nResumen: {len(slides)} lámina(s), NO superado.")
        sys.exit(1)
    print(f"✓ {len(slides)} lámina(s), cada una con una imagen a sangre. Validación superada.")
    sys.exit(0)


if __name__ == "__main__":
    main()
